from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
DOCS = ROOT / "docs"


SLIDES = [
    ("SecureLink", "Production-grade encrypted messaging and attack detection platform"),
    ("Snowflake / UNC5537", "The June 2024 customer data theft campaign showed how stolen credentials and missing MFA can lead to cloud data theft and extortion."),
    ("What Failed", "Stolen credentials, missing MFA, weak network allow lists, stale secrets, and limited anomaly detection created avoidable exposure."),
    ("Cryptography", "AES-GCM, SHA-256, HMAC, RSA-PSS signatures, Diffie-Hellman, and JWT authentication."),
    ("Architecture", "FastAPI, PostgreSQL, Streamlit, Prometheus, Grafana, Redis, and Caddy TLS."),
    ("Message Flow", "JWT auth, DH-derived session key, AES-GCM encryption, HMAC, RSA signature, verification before decrypt."),
    ("Attack Simulations", "Replay, tampering, invalid signature, MITM, and brute-force login attempts are blocked or logged."),
    ("Monitoring", "Security events appear in Streamlit; API metrics are exported to Prometheus and Grafana."),
    ("Recommendations", "Use TLS, rate limiting, key rotation, audit logging, least privilege, and continuous monitoring."),
    ("Conclusion", "SecureLink demonstrates layered cryptographic controls and intrusion-detection style visibility."),
]


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def box(draw: ImageDraw.ImageDraw, xy: tuple[int, int, int, int], title: str, body: str, outline: str) -> None:
    draw.rounded_rectangle(xy, radius=10, fill="white", outline=outline, width=3)
    draw.text((xy[0] + 25, xy[1] + 28), title, fill="#17202a", font=font(22, True))
    draw.text((xy[0] + 25, xy[1] + 58), body, fill="#53616f", font=font(17))


def arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int]) -> None:
    draw.line([start, end], fill="#53616f", width=4)
    draw.polygon([(end[0], end[1]), (end[0] - 14, end[1] - 8), (end[0] - 14, end[1] + 8)], fill="#53616f")


def canvas(title: str) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    image = Image.new("RGB", (1200, 680), "#f7f9fb")
    draw = ImageDraw.Draw(image)
    draw.text((60, 45), title, fill="#17202a", font=font(36, True))
    return image, draw


def render_architecture_png(path: Path) -> None:
    image, draw = canvas("SecureLink Architecture")
    box(draw, (60, 160, 300, 255), "Users", "Alice, Bob, demos", "#2f6fed")
    box(draw, (60, 335, 300, 430), "Attack Simulator", "Replay, tamper, MITM", "#d64545")
    draw.rounded_rectangle((405, 130, 720, 490), radius=10, fill="white", outline="#17202a", width=3)
    draw.text((455, 170), "FastAPI Backend", fill="#17202a", font=font(26, True))
    for i, item in enumerate(["Auth + JWT + bcrypt", "Messages API", "Crypto Services", "Detection + Audit"]):
        y = 220 + i * 65
        draw.rounded_rectangle((445, y, 680, y + 45), radius=8, fill="#eef5ff", outline="#2f6fed", width=2)
        draw.text((465, y + 12), item, fill="#17202a", font=font(17))
    box(draw, (835, 150, 1100, 245), "PostgreSQL", "Users, messages, events", "#1f8f5f")
    box(draw, (835, 315, 1100, 410), "Streamlit Dashboard", "Security event view", "#8d5cf6")
    box(draw, (835, 480, 1100, 575), "Prometheus + Grafana", "Metrics and monitoring", "#53616f")
    for start, end in [((300, 205), (405, 250)), ((300, 380), (405, 385)), ((720, 250), (835, 198)), ((720, 385), (835, 362)), ((720, 445), (835, 525))]:
        arrow(draw, start, end)
    image.save(path)


def render_attack_flow_png(path: Path) -> None:
    image, draw = canvas("Attack Flow and Defenses")
    box(draw, (70, 145, 320, 230), "Tampering", "Ciphertext changed", "#d64545")
    box(draw, (70, 275, 320, 360), "Replay", "Nonce or ID reused", "#d64545")
    box(draw, (70, 405, 320, 490), "Impersonation", "Invalid signature", "#d64545")
    draw.rounded_rectangle((465, 145, 740, 490), radius=10, fill="white", outline="#17202a", width=3)
    draw.text((510, 190), "Verification Gate", fill="#17202a", font=font(26, True))
    for i, line in enumerate(["Validate timestamp", "Reject duplicate nonce", "Verify HMAC", "Verify RSA signature", "Decrypt AES-GCM"]):
        draw.text((510, 245 + i * 45), f"{i + 1}. {line}", fill="#17202a", font=font(19))
    box(draw, (890, 220, 1130, 305), "Accepted", "Plaintext returned", "#1f8f5f")
    box(draw, (890, 370, 1130, 455), "Blocked + Logged", "Security event stored", "#d99000")
    for start, end in [((320, 188), (465, 255)), ((320, 318), (465, 318)), ((320, 448), (465, 390)), ((740, 280), (890, 262)), ((740, 390), (890, 410))]:
        arrow(draw, start, end)
    image.save(path)


def render_dashboard_png(path: Path) -> None:
    image, draw = canvas("SecureLink Security Dashboard")
    draw.rounded_rectangle((60, 130, 1140, 610), radius=10, fill="white", outline="#d5dce3", width=2)
    for x, title, value, outline in [(90, "Total Events", "8", "#2f6fed"), (450, "High Severity", "3", "#d64545"), (810, "Replay Attempts", "1", "#d99000")]:
        draw.rounded_rectangle((x, 170, x + 300, 285), radius=10, fill="#fbfcfd", outline=outline, width=3)
        draw.text((x + 25, 200), title, fill="#53616f", font=font(17))
        draw.text((x + 25, 230), value, fill="#17202a", font=font(42, True))
    draw.rounded_rectangle((90, 340, 1110, 560), radius=10, fill="#fbfcfd", outline="#d5dce3", width=2)
    draw.text((115, 375), "Recent Security Events", fill="#17202a", font=font(22, True))
    rows = [
        ("REPLAY_ATTACK_DETECTED", "HIGH", "Rejected duplicate message id or nonce"),
        ("HMAC_VERIFICATION_FAILED", "HIGH", "HMAC failed for tampered message"),
    ]
    for i, row in enumerate(rows):
        y = 430 + i * 45
        draw.text((115, y), row[0], fill="#17202a", font=font(17))
        draw.text((430, y), row[1], fill="#d64545", font=font(17, True))
        draw.text((610, y), row[2], fill="#17202a", font=font(17))
    image.save(path)


def export_png_assets() -> None:
    render_architecture_png(ASSETS / "architecture.png")
    render_attack_flow_png(ASSETS / "attack-flow.png")
    render_dashboard_png(ASSETS / "dashboard-preview.png")


def build_slide_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for title, body in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(12), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(34)
        title_frame.paragraphs[0].font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(11.6), Inches(1.8))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.text = body
        body_frame.paragraphs[0].font.size = Pt(22)

        if title == "Architecture":
            slide.shapes.add_picture(str(ASSETS / "architecture.png"), Inches(1.0), Inches(3.0), width=Inches(11.3))
        elif title == "Attack Simulations":
            slide.shapes.add_picture(str(ASSETS / "attack-flow.png"), Inches(1.0), Inches(3.0), width=Inches(11.3))
        elif title == "Monitoring":
            slide.shapes.add_picture(str(ASSETS / "dashboard-preview.png"), Inches(1.0), Inches(3.0), width=Inches(11.3))

    prs.save(DOCS / "SecureLink_presentation.pptx")


def main() -> None:
    export_png_assets()
    build_slide_deck()


if __name__ == "__main__":
    main()
